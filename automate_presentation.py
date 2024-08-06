import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from loguru import logger
from pathlib import Path
import sys


# Define custom colors for presentation
DARK_BLUE = RGBColor(0, 32, 96)
LIGHT_BLUE = RGBColor(197, 217, 241)


def load_and_clean_data(file_path: str) -> pd.DataFrame:
    """Loads the Netflix dataset, cleans it and preprocesses it.

    Args:
        file_path (str): The path to the CSV file.

    Returns:
        pd.DataFrame: The cleaned and processed DataFrame.
    """
    try:
        df = pd.read_csv(file_path)
        # Drop null values
        df.dropna(inplace=True)

        # Convert duration to minutes for movies
        df['duration_min'] = df[df['type'] == 'Movie']['duration'].str.extract('(\d+)').astype(float)

        # Add 'genre' column
        df['genre'] = df['listed_in'].str.split(',').str[0]

        # Convert 'date_added' to datetime
        df['year_added'] = pd.to_datetime(df['date_added'], format='mixed').dt.year
        
        logger.info("Data loaded and cleaned successfully.")
        return df
    except FileNotFoundError as e:
        logger.error(f"Error loading data: {e}")
        sys.exit(1)  # Exit with error code

def create_chart(data: pd.Series, title: str, filename: str, kind: str = 'bar', figsize: tuple = (10, 6)) -> None:
    """Creates and saves a chart.

    Args:
        data (pd.Series): The data to plot.
        title (str): The title of the chart.
        filename (str): The name of the file to save the chart as.
        kind (str, optional): The type of chart ('bar' or 'line'). Defaults to 'bar'.
        figsize (tuple, optional): The size of the figure. Defaults to (10, 6).
    """
    plt.figure(figsize=figsize)
    sns.set_style("whitegrid")

    if kind == 'bar':
        sns.barplot(x=data.index, y=data.values)
    elif kind == 'line':
        sns.lineplot(x=data.index, y=data.values)
    elif kind == 'hist':
        sns.histplot(data, kde=True)
    elif kind == 'scatter':
        sns.scatterplot(data=data, x='release_year', y='duration_min')
    elif kind == 'count':
        sns.countplot(data=data, y='rating', order=data['rating'].value_counts().index)

    plt.title(title, fontsize=16)
    plt.xlabel('')  # Remove default x-axis label
    plt.xticks(rotation=45)
    plt.tight_layout()

    img_dir = Path("img/slides")
    img_dir.mkdir(parents=True, exist_ok=True)

    plt.savefig(img_dir / filename)
    logger.info(f"Saved chart: {filename}")
    plt.close()


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """Adds a title slide to the presentation.

    Args:
        prs (Presentation): The presentation object.
        title (str): The title of the slide.
        subtitle (str): The subtitle of the slide.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)


def add_content_slide(prs: Presentation, title: str, image_path: str) -> None:
    """Adds a content slide with an image to the presentation.

    Args:
        prs (Presentation): The presentation object.
        title (str): The title of the slide.
        image_path (str): The path to the image file.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5.5)
    slide.shapes.add_picture(image_path, left, top, width, height)

    # Customize title font
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)

def customize_slide_master(prs: Presentation) -> None:
    """Customizes the slide master of the presentation.

    Args:
        prs (Presentation): The presentation object.
    """
    slide_master = prs.slide_master
    background = slide_master.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BLUE

    title_style = slide_master.slide_layouts[0].placeholders[0].text_frame.paragraphs[0].font
    title_style.name = 'Arial'
    title_style.size = Pt(44)
    title_style.color.rgb = DARK_BLUE

    body_style = slide_master.slide_layouts[1].placeholders[1].text_frame.paragraphs[0].font
    body_style.name = 'Arial'
    body_style.size = Pt(18)
    body_style.color.rgb = RGBColor(0, 0, 0)



def main():
    logger.info("Starting Netflix content analysis...")

    # Load and clean data
    df = load_and_clean_data('./data/netflix_titles.csv')

    # Generate Charts
    create_chart(df['type'].value_counts(), 'Movies vs. TV Shows', 'content_types.png')
    create_chart(df['country'].value_counts().nlargest(10), 'Top 10 Countries Producing Netflix Content', 'top_countries.png')
    create_chart(df['year_added'].value_counts().sort_index(), 'Content Added by Year', 'content_by_year.png', kind='line')
    create_chart(df['genre'].value_counts().nlargest(10), 'Top 10 Genres on Netflix', 'top_genres.png')

    # Movie-specific charts
    create_chart(df[df['type'] == 'Movie']['duration_min'], 'Distribution of Movie Durations', 'movie_duration_dist.png', kind='hist')
    create_chart(df[df['type'] == 'Movie'], 'Movie Duration vs. Release Year', 'duration_vs_year.png', kind='scatter')
    create_chart(df, 'Distribution of Content Ratings', 'rating_distribution.png', kind='count')


    # Create Presentation
    prs = Presentation()
    add_title_slide(prs, "Netflix Content Analysis", "Insights from the Netflix Movies and TV Shows Dataset")

    # Add slides with images
    img_dir = Path("img/slides")
    for img_path in img_dir.iterdir():
        add_content_slide(prs, img_path.stem.replace('_', ' ').title(), str(img_path))

    # Save the presentation
    prs.save('Netflix_Content_Analysis.pptx')
    logger.info("Presentation saved as Netflix_Content_Analysis.pptx")

    # Customize and Save Themed Presentation
    customize_slide_master(prs)
    prs.save('Netflix_Content_Analysis_Themed.pptx')
    logger.info("Themed presentation saved as Netflix_Content_Analysis_Themed.pptx")


if __name__ == "__main__":
    main()

